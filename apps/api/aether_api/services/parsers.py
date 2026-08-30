from __future__ import annotations

from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from io import BytesIO


@dataclass
class ParseResult:
    text: str
    pages: int = 0
    notices: list[str] = field(default_factory=list)  # e.g. "scanned PDF: OCR needed"


class DocumentParserProvider(ABC):
    @abstractmethod
    def can_parse(self, mime: str) -> bool: ...

    @abstractmethod
    def parse(self, data: bytes, mime: str, filename: str) -> ParseResult: ...


class TextParser(DocumentParserProvider):
    def can_parse(self, mime: str) -> bool:
        return mime.startswith("text/") or mime == "application/json"

    def parse(self, data: bytes, mime: str, filename: str) -> ParseResult:
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = data.decode("gbk")
            except UnicodeDecodeError:
                text = data.decode("utf-8", errors="replace")
        return ParseResult(text=text)


class PdfParser(DocumentParserProvider):
    def can_parse(self, mime: str) -> bool:
        return mime == "application/pdf"

    def parse(self, data: bytes, mime: str, filename: str) -> ParseResult:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        pages = []
        for page in reader.pages:
            try:
                pages.append(page.extract_text() or "")
            except Exception:  # noqa: BLE001
                pages.append("")
        text = "\n\n".join(p.strip() for p in pages).strip()
        notices = []
        if not text:
            notices.append("No extractable text layer (scanned PDF). OCR arrives in a later phase.")
        return ParseResult(text=text, pages=len(reader.pages), notices=notices)


class DocxParser(DocumentParserProvider):
    def can_parse(self, mime: str) -> bool:
        return mime.endswith("wordprocessingml.document")

    def parse(self, data: bytes, mime: str, filename: str) -> ParseResult:
        import docx

        doc = docx.Document(BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        return ParseResult(text="\n".join(parts))


class XlsxParser(DocumentParserProvider):
    def can_parse(self, mime: str) -> bool:
        return mime.endswith("spreadsheetml.sheet")

    def parse(self, data: bytes, mime: str, filename: str) -> ParseResult:
        import openpyxl

        wb = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
        wb.close()
        return ParseResult(text="\n".join(parts))


class PptxParser(DocumentParserProvider):
    def can_parse(self, mime: str) -> bool:
        return mime.endswith("presentationml.presentation")

    def parse(self, data: bytes, mime: str, filename: str) -> ParseResult:
        from pptx import Presentation

        prs = Presentation(BytesIO(data))
        parts = []
        slides = list(prs.slides)
        for i, slide in enumerate(slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return ParseResult(text="\n".join(parts), pages=len(slides))


_PARSERS: list[DocumentParserProvider] = [
    PdfParser(), DocxParser(), XlsxParser(), PptxParser(), TextParser(),
]


def register_parser(parser: DocumentParserProvider) -> None:
    _PARSERS.insert(0, parser)


def get_parser(mime: str) -> DocumentParserProvider | None:
    for parser in _PARSERS:
        if parser.can_parse(mime):
            return parser
    return None
